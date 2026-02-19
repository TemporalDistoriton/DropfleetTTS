"""
Upgrade Card Generator
======================
Reads UpgradeTemplate.pptx and Upgrades.xlsx from the same folder as this script.
For each row in the Excel sheet, replaces placeholders in slide 1, exports as PNG,
crops the white PowerPoint canvas margin, applies a faction colour tint, saves to ./output/

Excel columns required: Name, Effect, Ships, Points, Faction

Faction values (case-insensitive): phr, scourge, resistance, shaltari, ucm, bioficers, civilian

Dependencies:
    pip install python-pptx openpyxl Pillow numpy

On Windows, PowerPoint must be installed (used via COM automation for PNG export).
    pip install comtypes   (Windows only)

On Linux/Mac, LibreOffice (soffice) and Poppler (pdftoppm) must be on PATH.

Usage:
    python generate_upgrade_cards.py
"""

import os
import sys
import platform
import subprocess
import tempfile

try:
    from pptx import Presentation
except ImportError:
    sys.exit("Missing: pip install python-pptx")

try:
    import openpyxl
except ImportError:
    sys.exit("Missing: pip install openpyxl")

try:
    from PIL import Image
    import numpy as np
except ImportError:
    sys.exit("Missing: pip install Pillow numpy")

# ── faction tint colours (RGB 0-1 floats) ─────────────────────────────────────
FACTION_TINTS = {
    "phr":        (0.58,  0.498, 0.231),
    "scourge":    (0.431, 0.231, 0.58),
    "resistance": (0.231, 0.463, 0.58),
    "shaltari":   (0.58,  0.325, 0.231),
    "ucm":        (0.341, 0.569, 0.337),
    "bioficers":  (0.58,  0.231, 0.231),
    "civilian":   (0.231, 0.576, 0.58),
}

# Placeholder strings exactly as they appear in UpgradeTemplate.pptx
PH_NAME   = "UPGRADE NAME"
PH_EFFECT = "EFFECT"
PH_SHIPS  = "SHIPS ABLE TO take this upgrade"
PH_POINTS = "1"

TINT_STRENGTH = 0.20
RENDER_DPI    = 150


def replace_text_in_shape(shape, old: str, new: str):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def export_slide_png_windows(pptx_path: str, png_path: str):
    """Export slide 1 to PNG using installed PowerPoint via COM (Windows only)."""
    try:
        import comtypes.client
    except ImportError:
        sys.exit(
            "\nMissing Windows dependency: pip install comtypes\n"
            "This is needed to drive PowerPoint for PNG export on Windows.\n"
        )

    abs_pptx = os.path.abspath(pptx_path)
    abs_png  = os.path.abspath(png_path)

    powerpoint = None
    prs_com    = None
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1

        prs_com = powerpoint.Presentations.Open(
            abs_pptx,
            ReadOnly=True,
            Untitled=False,
            WithWindow=False,
        )

        slide   = prs_com.Slides(1)
        slide_w = prs_com.PageSetup.SlideWidth
        slide_h = prs_com.PageSetup.SlideHeight
        export_w = 2000
        export_h = int(export_w * slide_h / slide_w)

        slide.Export(abs_png, "PNG", export_w, export_h)

    finally:
        if prs_com is not None:
            try:
                prs_com.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass


def export_slide_png_libreoffice(pptx_path: str, png_path: str, dpi: int = 150):
    """Export slide 1 to PNG using LibreOffice + pdftoppm (Linux/Mac)."""
    import shutil as _shutil
    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", tmpdir, pptx_path],
            capture_output=True, text=True
        )
        pdf_files = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]
        if not pdf_files:
            raise RuntimeError(
                f"LibreOffice failed.\nSTDOUT: {result.stdout}\nSTDERR: {result.stderr}"
            )
        pdf_path = os.path.join(tmpdir, pdf_files[0])
        base     = os.path.join(tmpdir, "slide")
        result2  = subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), "-f", "1", "-l", "1", pdf_path, base],
            capture_output=True, text=True
        )
        png_files = sorted(f for f in os.listdir(tmpdir) if f.endswith(".png"))
        if not png_files:
            raise RuntimeError(
                f"pdftoppm failed.\nSTDOUT: {result2.stdout}\nSTDERR: {result2.stderr}"
            )
        _shutil.copy(os.path.join(tmpdir, png_files[0]), png_path)


def export_slide_to_png(pptx_path: str, png_path: str):
    if platform.system() == "Windows":
        export_slide_png_windows(pptx_path, png_path)
    else:
        export_slide_png_libreoffice(pptx_path, png_path, dpi=RENDER_DPI)


def crop_white_border(img: Image.Image, threshold: int = 240) -> Image.Image:
    img_rgb = img.convert("RGB")
    arr     = np.array(img_rgb)
    mask    = ~np.all(arr >= threshold, axis=2)
    rows    = np.any(mask, axis=1)
    cols    = np.any(mask, axis=0)
    if not rows.any():
        return img
    rmin = max(0,              int(np.where(rows)[0][0])  - 2)
    rmax = min(arr.shape[0]-1, int(np.where(rows)[0][-1]) + 2)
    cmin = max(0,              int(np.where(cols)[0][0])  - 2)
    cmax = min(arr.shape[1]-1, int(np.where(cols)[0][-1]) + 2)
    return img.crop((cmin, rmin, cmax + 1, rmax + 1))


def apply_tint(img: Image.Image, tint_rgb_01: tuple, strength: float = 0.20) -> Image.Image:
    img  = img.convert("RGBA")
    tr, tg, tb = [int(c * 255) for c in tint_rgb_01]
    arr  = np.array(img, dtype=np.float32)
    arr[:, :, 0] = arr[:, :, 0] * (1 - strength) + tr * strength
    arr[:, :, 1] = arr[:, :, 1] * (1 - strength) + tg * strength
    arr[:, :, 2] = arr[:, :, 2] * (1 - strength) + tb * strength
    return Image.fromarray(np.clip(arr, 0, 255).astype(np.uint8), "RGBA")


def make_corners_transparent(img: Image.Image, tolerance: int = 30) -> Image.Image:
    """
    Flood-fill from all four corners to make the background triangles transparent.
    Works by replacing the corner colour (and nearby colours within *tolerance*)
    with full transparency, leaving the octagonal card shape intact.
    """
    from PIL import ImageDraw
    img  = img.convert("RGBA")
    w, h = img.size
    for corner in [(0, 0), (w - 1, 0), (0, h - 1), (w - 1, h - 1)]:
        ImageDraw.floodfill(img, corner, (0, 0, 0, 0), thresh=tolerance)
    return img


def sanitise_filename(name: str) -> str:
    keep = set(" abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789_-.")
    return "".join(c if c in keep else "_" for c in name).strip()


def main():
    script_dir    = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "UpgradeTemplate.pptx")
    excel_path    = os.path.join(script_dir, "Upgrades.xlsx")
    output_dir    = os.path.join(script_dir, "output")

    for p in (template_path, excel_path):
        if not os.path.isfile(p):
            sys.exit(f"ERROR: Cannot find '{p}'")

    os.makedirs(output_dir, exist_ok=True)

    wb      = openpyxl.load_workbook(excel_path, data_only=True)
    ws      = wb.active
    headers = [str(c.value).strip() if c.value else "" for c in next(ws.iter_rows(min_row=1, max_row=1))]

    def col(row_vals, col_name: str) -> str:
        try:
            v = row_vals[headers.index(col_name)]
            return str(v) if v is not None else ""
        except ValueError:
            return ""

    rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if any(r)]
    print(f"Found {len(rows)} upgrade(s) to process.\n")

    for row_vals in rows:
        name    = col(row_vals, "Name")
        effect  = col(row_vals, "Effect")
        ships   = col(row_vals, "Ships")
        points  = col(row_vals, "Points")
        faction = col(row_vals, "Faction").strip().lower()

        if not name:
            print("  Skipping row with empty Name.")
            continue

        print(f"  Processing: {name}  [{faction}]")

        prs   = Presentation(template_path)
        slide = prs.slides[0]

        for shape in slide.shapes:
            replace_text_in_shape(shape, PH_NAME,   name)
            replace_text_in_shape(shape, PH_EFFECT, effect)
            replace_text_in_shape(shape, PH_SHIPS,  ships)
            replace_text_in_shape(shape, PH_POINTS, points)

        while len(prs.slides) > 1:
            sldIdLst = prs.slides._sldIdLst
            last     = sldIdLst[-1]
            rId      = last.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            sldIdLst.remove(last)
            if rId and rId in prs.part.related_parts:
                del prs.part.related_parts[rId]

        tmp_fd, tmp_pptx = tempfile.mkstemp(suffix=".pptx")
        os.close(tmp_fd)

        try:
            prs.save(tmp_pptx)
            raw_png = os.path.join(output_dir, f"_raw_{sanitise_filename(name)}.png")
            export_slide_to_png(tmp_pptx, raw_png)
        finally:
            if os.path.exists(tmp_pptx):
                os.unlink(tmp_pptx)

        img = Image.open(raw_png).convert("RGBA")
        img = crop_white_border(img)

        tint = FACTION_TINTS.get(faction)
        if tint:
            img = apply_tint(img, tint, strength=TINT_STRENGTH)
        else:
            print(f"    WARNING: Unknown faction '{faction}' – no tint applied.")

        # Make octagonal corners transparent
        img = make_corners_transparent(img)

        out_path = os.path.join(output_dir, f"{sanitise_filename(name)}.png")
        img.save(out_path, "PNG")
        os.unlink(raw_png)
        print(f"    Saved -> {out_path}")

    print("\nAll done!")


if __name__ == "__main__":
    main()